import os
import pickle
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import argparse
from datetime import datetime

# Document processing
import PyPDF2
import docx
from pptx import Presentation

# Text processing and embeddings
from sentence_transformers import SentenceTransformer
import nltk
from nltk.tokenize import sent_tokenize
import re

# Vector store
import faiss
import numpy as np

# LLM integration
import httpx
from dotenv import load_dotenv

# Optional: LangChain integration
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain.schema import Document
    LANGCHAIN_AVAILABLE = True
except ImportError:
    LANGCHAIN_AVAILABLE = False
    print("LangChain не установлен. Используется базовый text splitter.")

# Load environment variables
load_dotenv()

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Класс для обработки различных типов документов"""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Извлекает текст из PDF файла"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            logger.error(f"Ошибка при чтении PDF {file_path}: {e}")
        return text
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Извлекает текст из Word документа"""
        text = ""
        try:
            doc = docx.Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            logger.error(f"Ошибка при чтении DOCX {file_path}: {e}")
        return text
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Извлекает текст из PowerPoint презентации"""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            logger.error(f"Ошибка при чтении PPTX {file_path}: {e}")
        return text
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Извлекает текст из текстового файла"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='cp1251') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Ошибка при чтении TXT {file_path}: {e}")
        except Exception as e:
            logger.error(f"Ошибка при чтении TXT {file_path}: {e}")
        return ""
    
    @classmethod
    def extract_text(cls, file_path: str) -> str:
        """Извлекает текст из файла в зависимости от его типа"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        extractors = {
            '.pdf': cls.extract_text_from_pdf,
            '.docx': cls.extract_text_from_docx,
            '.doc': cls.extract_text_from_docx,
            '.pptx': cls.extract_text_from_pptx,
            '.ppt': cls.extract_text_from_pptx,
            '.txt': cls.extract_text_from_txt,
            '.md': cls.extract_text_from_txt,
        }
        
        if extension in extractors:
            return extractors[extension](str(file_path))
        else:
            logger.warning(f"Неподдерживаемый тип файла: {extension}")
            return ""

class TextSplitter:
    """Класс для разбиения текста на чанки"""
    
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50, use_langchain: bool = True):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.use_langchain = use_langchain and LANGCHAIN_AVAILABLE
        
        if self.use_langchain:
            self.splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
            )
        
        # Download NLTK data if needed
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt', quiet=True)
    
    def split_text_langchain(self, text: str) -> List[str]:
        """Использует LangChain для разбиения текста"""
        docs = self.splitter.split_text(text)
        return docs
    
    def split_text_basic(self, text: str) -> List[str]:
        """Базовое разбиение текста на чанки"""
        # Разбиваем на предложения
        sentences = sent_tokenize(text)
        
        chunks = []
        current_chunk = ""
        current_length = 0
        
        for sentence in sentences:
            sentence_length = len(sentence)
            
            if current_length + sentence_length <= self.chunk_size:
                current_chunk += sentence + " "
                current_length += sentence_length
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                
                # Начинаем новый чанк
                if sentence_length <= self.chunk_size:
                    current_chunk = sentence + " "
                    current_length = sentence_length
                else:
                    # Разбиваем очень длинное предложение
                    words = sentence.split()
                    temp_chunk = ""
                    for word in words:
                        if len(temp_chunk) + len(word) <= self.chunk_size:
                            temp_chunk += word + " "
                        else:
                            if temp_chunk:
                                chunks.append(temp_chunk.strip())
                            temp_chunk = word + " "
                    current_chunk = temp_chunk
                    current_length = len(temp_chunk)
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return [chunk for chunk in chunks if len(chunk.strip()) > 20]
    
    def split_text(self, text: str) -> List[str]:
        """Разбивает текст на чанки"""
        if self.use_langchain:
            return self.split_text_langchain(text)
        else:
            return self.split_text_basic(text)

class VectorStore:
    """Класс для работы с векторным хранилищем FAISS"""
    
    def __init__(self, embedding_model: str = "all-MiniLM-L6-v2"):
        self.embedding_model_name = embedding_model
        self.model = SentenceTransformer(embedding_model)
        self.dimension = self.model.get_sentence_embedding_dimension()
        
        # Инициализируем FAISS индекс
        self.index = faiss.IndexFlatIP(self.dimension)  # Inner Product для косинусного сходства
        self.documents = []  # Хранилище исходных текстов
        self.metadata = []   # Метаданные документов
        
        logger.info(f"Инициализирован FAISS индекс с размерностью {self.dimension}")
    
    def add_documents(self, texts: List[str], metadatas: Optional[List[Dict]] = None):
        """Добавляет документы в векторное хранилище"""
        if metadatas is None:
            metadatas = [{} for _ in texts]
        
        logger.info(f"Добавляем {len(texts)} документов в векторное хранилище...")
        
        # Создаем эмбеддинги
        embeddings = self.model.encode(texts, convert_to_tensor=False, show_progress_bar=True)
        
        # Нормализуем для косинусного сходства
        embeddings = embeddings / np.linalg.norm(embeddings, axis=1, keepdims=True)
        
        # Добавляем в FAISS
        self.index.add(embeddings.astype('float32'))
        
        # Сохраняем тексты и метаданные
        self.documents.extend(texts)
        self.metadata.extend(metadatas)
        
        logger.info(f"Добавлено {len(texts)} документов. Всего в индексе: {len(self.documents)}")
    
    def search(self, query: str, k: int = 5) -> List[Dict[str, Any]]:
        """Поиск наиболее релевантных документов"""
        if len(self.documents) == 0:
            return []
        
        # Создаем эмбеддинг запроса
        query_embedding = self.model.encode([query])
        query_embedding = query_embedding / np.linalg.norm(query_embedding, axis=1, keepdims=True)
        
        # Поиск в FAISS
        scores, indices = self.index.search(query_embedding.astype('float32'), k)
        
        results = []
        for score, idx in zip(scores[0], indices[0]):
            if idx < len(self.documents):
                results.append({
                    'content': self.documents[idx],
                    'metadata': self.metadata[idx],
                    'score': float(score)
                })
        
        return results
    
    def save(self, path: str):
        """Сохраняет векторное хранилище"""
        save_data = {
            'documents': self.documents,
            'metadata': self.metadata,
            'embedding_model_name': self.embedding_model_name
        }
        
        # Сохраняем FAISS индекс
        faiss.write_index(self.index, f"{path}.faiss")
        
        # Сохраняем остальные данные
        with open(f"{path}.pkl", 'wb') as f:
            pickle.dump(save_data, f)
        
        logger.info(f"Векторное хранилище сохранено: {path}")
    
    def load(self, path: str):
        """Загружает векторное хранилище"""
        # Загружаем FAISS индекс
        self.index = faiss.read_index(f"{path}.faiss")
        
        # Загружаем остальные данные
        with open(f"{path}.pkl", 'rb') as f:
            save_data = pickle.load(f)
        
        self.documents = save_data['documents']
        self.metadata = save_data['metadata']
        
        # Проверяем совместимость модели
        if save_data['embedding_model_name'] != self.embedding_model_name:
            logger.warning(f"Модель эмбеддингов изменилась: {save_data['embedding_model_name']} -> {self.embedding_model_name}")
        
        logger.info(f"Векторное хранилище загружено: {len(self.documents)} документов")

class LLMClient:
    """Клиент для работы с LLM"""
    
    def __init__(self, provider: str = "mistral", model_name: str = "mistral-small"):
        self.provider = provider
        self.model_name = model_name
        
        if provider == "mistral":
            self.api_key = os.getenv("MISTRAL_API_KEY")
            self.base_url = "https://api.mistral.ai/v1"
        elif provider == "local":
            self.base_url = os.getenv("LOCAL_LLM_URL", "http://localhost:8000")
            self.api_key = None
        else:
            raise ValueError(f"Неподдерживаемый провайдер: {provider}")
    
    async def generate_response(self, prompt: str, max_tokens: int = 1000) -> str:
        """Генерирует ответ от LLM"""
        async with httpx.AsyncClient() as client:
            if self.provider == "mistral":
                headers = {
                    "Authorization": f"Bearer {self.api_key}",
                    "Content-Type": "application/json"
                }
                
                payload = {
                    "model": self.model_name,
                    "messages": [{"role": "user", "content": prompt}],
                    "max_tokens": max_tokens,
                    "temperature": 0.7
                }
                
                response = await client.post(
                    f"{self.base_url}/chat/completions",
                    json=payload,
                    headers=headers,
                    timeout=30.0
                )
                response.raise_for_status()
                
                data = response.json()
                return data["choices"][0]["message"]["content"]
            
            elif self.provider == "local":
                payload = {
                    "blocks": [
                        {"trigger": "user", "message": prompt}
                    ],
                    "max_tokens": max_tokens,
                    "temperature": 0.7
                }
                
                response = await client.post(
                    f"{self.base_url}/generate",
                    json=payload,
                    timeout=30.0
                )
                response.raise_for_status()
                
                data = response.json()
                return data["response"]

class RAGChatbot:
    """Основной класс RAG чатбота"""
    
    def __init__(self, 
                 embedding_model: str = "all-MiniLM-L6-v2",
                 llm_provider: str = "mistral",
                 llm_model: str = "mistral-small",
                 chunk_size: int = 500,
                 chunk_overlap: int = 50):
        
        self.vector_store = VectorStore(embedding_model)
        self.text_splitter = TextSplitter(chunk_size, chunk_overlap)
        self.llm_client = LLMClient(llm_provider, llm_model)
        self.document_processor = DocumentProcessor()
        
        logger.info(f"RAG чатбот инициализирован:")
        logger.info(f"  - Embedding model: {embedding_model}")
        logger.info(f"  - LLM: {llm_provider}/{llm_model}")
        logger.info(f"  - Chunk size: {chunk_size}")
    
    def load_documents(self, file_paths: List[str]) -> int:
        """Загружает и индексирует документы"""
        all_chunks = []
        all_metadata = []
        
        for file_path in file_paths:
            logger.info(f"Обрабатываем файл: {file_path}")
            
            # Извлекаем текст
            text = self.document_processor.extract_text(file_path)
            
            if not text.strip():
                logger.warning(f"Файл {file_path} пустой или не удалось извлечь текст")
                continue
            
            # Разбиваем на чанки
            chunks = self.text_splitter.split_text(text)
            logger.info(f"Создано {len(chunks)} чанков из {file_path}")
            
            # Создаем метаданные
            metadata = [{
                'source': file_path,
                'chunk_id': i,
                'total_chunks': len(chunks),
                'file_type': Path(file_path).suffix
            } for i in range(len(chunks))]
            
            all_chunks.extend(chunks)
            all_metadata.extend(metadata)
        
        # Добавляем в векторное хранилище
        if all_chunks:
            self.vector_store.add_documents(all_chunks, all_metadata)
        
        return len(all_chunks)
    
    def create_prompt(self, query: str, context_docs: List[Dict[str, Any]]) -> str:
        """Создает промпт для LLM с контекстом"""
        context = "\n\n".join([
            f"Документ {i+1} (score: {doc['score']:.3f}):\n{doc['content']}"
            for i, doc in enumerate(context_docs)
        ])
        
        prompt = f"""Ты полезный ассистент, который отвечает на вопросы на основе предоставленного контекста.

КОНТЕКСТ:
{context}

ВОПРОС: {query}

ИНСТРУКЦИИ:
1. Отвечай только на основе предоставленного контекста
2. Если информации недостаточно, честно скажи об этом
3. Указывай источники информации когда это возможно
4. Будь конкретным и полезным

ОТВЕТ:"""
        
        return prompt
    
    async def ask(self, query: str, k: int = 5) -> Dict[str, Any]:
        """Отвечает на вопрос пользователя"""
        if len(self.vector_store.documents) == 0:
            return {
                'answer': 'База знаний пуста. Загрузите документы для получения ответов.',
                'sources': [],
                'query': query
            }
        
        # Поиск релевантных документов
        logger.info(f"Поиск релевантных документов для запроса: {query}")
        relevant_docs = self.vector_store.search(query, k)
        
        if not relevant_docs:
            return {
                'answer': 'Не найдено релевантных документов для вашего вопроса.',
                'sources': [],
                'query': query
            }
        
        # Создаем промпт
        prompt = self.create_prompt(query, relevant_docs)
        
        # Генерируем ответ
        logger.info("Генерируем ответ от LLM...")
        answer = await self.llm_client.generate_response(prompt)
        
        # Подготавливаем источники
        sources = []
        for doc in relevant_docs:
            sources.append({
                'content': doc['content'][:200] + '...' if len(doc['content']) > 200 else doc['content'],
                'source': doc['metadata'].get('source', 'Unknown'),
                'score': doc['score']
            })
        
        return {
            'answer': answer,
            'sources': sources,
            'query': query,
            'timestamp': datetime.now().isoformat()
        }
    
    def save_knowledge_base(self, path: str):
        """Сохраняет базу знаний"""
        self.vector_store.save(path)
    
    def load_knowledge_base(self, path: str):
        """Загружает базу знаний"""
        self.vector_store.load(path)

async def interactive_chat(chatbot: RAGChatbot):
    """Интерактивный чат с пользователем"""
    print("\n" + "="*60)
    print("RAG Чатбот готов к общению!")
    print("Введите 'quit' для выхода")
    print("="*60 + "\n")
    
    while True:
        query = input("Вы: ").strip()
        
        if query.lower() in ['quit', 'exit', 'bye', 'выход']:
            print("До свидания!")
            break
        
        if not query:
            continue
        
        print("Думаю...")
        
        try:
            response = await chatbot.ask(query)
            
            print(f"\nБот: {response['answer']}\n")
            
            if response['sources']:
                print("Источники:")
                for i, source in enumerate(response['sources'][:3], 1):
                    print(f"  {i}. {source['source']} (relevance: {source['score']:.3f})")
                    print(f"     {source['content']}\n")
            
        except Exception as e:
            print(f"Ошибка: {e}")
            logger.error(f"Ошибка при обработке запроса: {e}")

def main():
    parser = argparse.ArgumentParser(description='RAG Chatbot Demo')
    parser.add_argument('--documents', '-d', nargs='+', help='Путь к документам для загрузки')
    parser.add_argument('--save-kb', help='Сохранить базу знаний в файл')
    parser.add_argument('--load-kb', help='Загрузить базу знаний из файла')
    parser.add_argument('--embedding-model', default='all-MiniLM-L6-v2', help='Модель для эмбеддингов')
    parser.add_argument('--llm-provider', choices=['mistral', 'local'], default='mistral', help='Провайдер LLM')
    parser.add_argument('--llm-model', default='mistral-small', help='Модель LLM')
    parser.add_argument('--chunk-size', type=int, default=500, help='Размер чанка текста')
    parser.add_argument('--interactive', action='store_true', help='Запустить в интерактивном режиме')
    parser.add_argument('--query', help='Задать один вопрос и получить ответ')
    
    args = parser.parse_args()
    
    # Создаем чатбот
    chatbot = RAGChatbot(
        embedding_model=args.embedding_model,
        llm_provider=args.llm_provider,
        llm_model=args.llm_model,
        chunk_size=args.chunk_size
    )
    
    # Загружаем существующую базу знаний
    if args.load_kb:
        try:
            chatbot.load_knowledge_base(args.load_kb)
            print(f"База знаний загружена из {args.load_kb}")
        except Exception as e:
            logger.error(f"Ошибка загрузки базы знаний: {e}")
    
    # Загружаем документы
    if args.documents:
        try:
            num_chunks = chatbot.load_documents(args.documents)
            print(f"Загружено {num_chunks} чанков из {len(args.documents)} документов")
        except Exception as e:
            logger.error(f"Ошибка загрузки документов: {e}")
            return
    
    # Сохраняем базу знаний
    if args.save_kb:
        try:
            chatbot.save_knowledge_base(args.save_kb)
            print(f"База знаний сохранена в {args.save_kb}")
        except Exception as e:
            logger.error(f"Ошибка сохранения базы знаний: {e}")
    
    # Обрабатываем запрос
    if args.query:
        import asyncio
        
        async def single_query():
            response = await chatbot.ask(args.query)
            print(f"\nВопрос: {response['query']}")
            print(f"Ответ: {response['answer']}\n")
            
            if response['sources']:
                print("Источники:")
                for i, source in enumerate(response['sources'], 1):
                    print(f"  {i}. {source['source']} (relevance: {source['score']:.3f})")
        
        asyncio.run(single_query())
    
    # Интерактивный режим
    elif args.interactive:
        import asyncio
        asyncio.run(interactive_chat(chatbot))
    
    else:
        print("Используйте --interactive для интерактивного режима или --query для разового вопроса")
        print("Помощь: python rag_chatbot.py --help")

if __name__ == "__main__":
    main()